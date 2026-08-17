import sys
from pptx import Presentation

def extract_pptx_info(path):
    try:
        prs = Presentation(path)
        print(f"Opening: {path}")
        print(f"Total Slides: {len(prs.slides)}")
        print("-" * 30)
        
        for i, slide in enumerate(prs.slides):
            print(f"Slide {i+1}:")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    print(f"  [Text]: {shape.text.strip()}")
            print("-" * 20)
    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    pptx_path = r"c:\Users\premg\OneDrive\Desktop\Major project\Project PPT -FOOG GUARD.pptx"
    extract_pptx_info(pptx_path)
