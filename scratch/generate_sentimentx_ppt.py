from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def create_sentimentx_ppt():
    prs = Presentation()

    # Define some styling constants (academic/clean)
    def set_slide_title(slide, text):
        title_shape = slide.shapes.title
        title_shape.text = text
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(36)
            paragraph.font.bold = True

    def add_bullet_points(slide, points):
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        for point in points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(18)

    # --- Slide 1: Title Slide ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "SentimentX: A Multilingual Hybrid CNN-XLNet Framework for Sentiment Analysis"
    
    subtitle = slide.placeholders[1]
    subtitle.text = (
        "Project Topic: Advanced Deep Learning in NLP\n"
        "PRESENTED BY:\n"
        "STUDENT NAME 1 (USN 1)\n"
        "STUDENT NAME 2 (USN 2)\n"
        "UNDER THE GUIDANCE:\n"
        "Mrs. DEEKSHA K R\n"
        "Department of Computer Science & Engineering\n"
        "Yenepoya Institute of Technology, Moodbidri"
    )

    # --- Slide 2: Table of Contents ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_title(slide, "TABLE OF CONTENTS")
    add_bullet_points(slide, [
        "ABSTRACT",
        "INTRODUCTION",
        "EXISTING SYSTEM",
        "PROPOSED SYSTEM",
        "SYSTEM REQUIREMENTS (HARDWARE & SOFTWARE)",
        "METHODOLOGY",
        "EXPECTED OUTCOME",
        "REFERENCES"
    ])

    # --- Slide 3: Abstract ---
    slide = prs.slides.add_slide(slide_layout)
    set_slide_title(slide, "ABSTRACT")
    add_bullet_points(slide, [
        "Sentiment analysis is crucial for understanding public opinion, yet most models are optimized exclusively for English.",
        "SentimentX addresses the gap in low-resource Indian languages (Kannada, Malayalam, Hindi, Marathi, Telugu).",
        "This project proposes a unique Hybrid CNN-XLNet architecture that leverages CNNs for spatial feature extraction and XLNet for deep semantic context.",
        "Tested on a verified dataset of 72,000+ records, the system provides high-precision emotion detection across 6 languages.",
        "The framework is designed for real-time application via a scalable FastAPI-driven dashboard."
    ])

    # --- Slide 4: Introduction ---
    slide = prs.slides.add_slide(slide_layout)
    set_slide_title(slide, "INTRODUCTION")
    add_bullet_points(slide, [
        "Language Diversity: India has 22+ official languages, yet digital sentiment tools often fail to capture regional nuances.",
        "Complex Script & Grammar: Indian languages present unique challenges like morphologically rich structures and complex tokenization.",
        "Code-Mixing: Modern social media uses a mix of native scripts and English alphabet (e.g., Tanglish, Hinglish).",
        "SentimentX bridges this divide by providing a unified transformer-based model optimized for these specific linguistic patterns."
    ])

    # --- Slide 5: Existing System ---
    slide = prs.slides.add_slide(slide_layout)
    set_slide_title(slide, "EXISTING SYSTEM")
    add_bullet_points(slide, [
        "Monolingual Bias: Most prevalent systems (BERT, VADER) are primarily English-centric.",
        "Poor Generalization: Existing models struggle when faced with code-mixed text or slang common in Indian dialects.",
        "Resource Scarcity: Lack of standardized, cleaned datasets for regional languages like Malayalam and Marathi.",
        "Speed Issues: Traditional RNN/LSTM based models are slow for real-time multilingual processing.",
        "Outcome: Leads to incorrect customer insights and misinterpretation of regional public sentiment."
    ])

    # --- Slide 6: Proposed System ---
    slide = prs.slides.add_slide(slide_layout)
    set_slide_title(slide, "PROPOSED SYSTEM")
    add_bullet_points(slide, [
        "SentimentX: A Multilingual Hybrid model combining Convolutional Neural Networks (CNN) and XLNet Transformers.",
        "Unprecedented Scale: Trained on a verified dataset of 72,365 records from 6 major archives.",
        "Emotion-Centric: Mapped to the Ekman 6-class schema (Happiness, Sadness, Anger, Fear, Surprise, Disgust).",
        "GPU Accelerated: Leverages NVIDIA CUDA for high-throughput training and inference.",
        "Unified API: A single endpoint that handles diverse Indian scripts and English interchangeably."
    ])

    # --- Slide 7: System Requirements ---
    slide = prs.slides.add_slide(slide_layout)
    set_slide_title(slide, "SYSTEM REQUIREMENTS")
    add_bullet_points(slide, [
        "Hardware Requirements:",
        "  • GPU: NVIDIA GeForce RTX 3050 (6GB VRAM)",
        "  • RAM: 16GB Dual-Channel",
        "  • Processor: Intel Core i5/i7 (11th Gen+)",
        "Software Requirements:",
        "  • Programming: Python 3.12",
        "  • Deep Learning: PyTorch, Transformers (Hugging Face)",
        "  • Backend & UI: FastAPI, Uvicorn, HTML5, CSS3, JavaScript"
    ])

    # --- Slide 8: Methodology ---
    slide = prs.slides.add_slide(slide_layout)
    set_slide_title(slide, "METHODOLOGY")
    add_bullet_points(slide, [
        "Data Merging: Aggregating 6 diverse Archives into language-specific clean CSVs.",
        "Hybrid Architecture: Integrating CNN layers as feature extractors for local patterns before passing to XLNet's global attention mechanism.",
        "Preprocessing: Advanced tokenization using SentencePiece to handle complex scripts.",
        "Optimized Training: Utilizing Class Weights for dataset balancing and Early Stopping to prevent overfitting.",
        "Validation: High-speed validation on a dedicated 20% holdout set (14k+ rows)."
    ])

    # --- Slide 9: Expected Outcome ---
    slide = prs.slides.add_slide(slide_layout)
    set_slide_title(slide, "EXPECTED OUTCOME & RESULTS")
    add_bullet_points(slide, [
        "High Performance: Achieved a Weighted Precision of 64.68% and F1-Score of 53.12%.",
        "Interactive Dashboard: A sleek, real-time web interface for sentiment analysis and prediction history.",
        "Extensible Framework: Easily handles new languages (like Malayalam/Telugu) with minimal retraining.",
        "Scalable API: Ready for integration with business chatbots or social media monitoring tools.",
        "Better Insights: Provides localized emotional context for regional Indian consumers."
    ])

    # --- Slide 10: References ---
    slide = prs.slides.add_slide(slide_layout)
    set_slide_title(slide, "REFERENCES")
    add_bullet_points(slide, [
        "[1] Yang, Z., et al. (2019). 'XLNet: Generalized Autoregressive Pretraining for Language Understanding.'",
        "[2] Kim, Y. (2014). 'Convolutional Neural Networks for Sentence Classification.'",
        "[3] Ekman, P. (1992). 'An Argument for Basic Emotions.' Cognition & Emotion.",
        "[4] SentimentX Repository Documentation and Internal Benchmarks (2026).",
        "[5] Hugging Face Transformers: State-of-the-Art Natural Language Processing."
    ])

    # --- Slide 11: Thank You ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "THANK YOU..."
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Q&A Session\nEmail: research@sentimentx.ai"

    # Save the presentation
    save_path = "SentimentX_Project_Presentation.pptx"
    prs.save(save_path)
    print(f"Presentation saved successfully at: {os.path.abspath(save_path)}")

if __name__ == "__main__":
    create_sentimentx_ppt()
